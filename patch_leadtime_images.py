# -*- coding: utf-8 -*-
"""Surgically swap ONLY the two lead-time figures (neupert.png on slide 3,
dashboard.png on slide 6) inside the already-built deck, in place.

Why this exists: the live .pptx contains manual PowerPoint edits that are NOT
in build_deck.py, so re-running the full build would destroy them. This script
opens the existing deck, finds each target picture by its on-slide position,
overwrites just that image part's bytes, and saves. No text, layout, or other
shape is touched.
"""
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

OUT = "SuryaSetu - ISRO BAH 2026 - PS15 - Idea Submission.pptx"
A = "render_assets"

# (slide index, new png, expected left/top in inches as placed by build_deck.py)
TARGETS = [
    (2, "neupert.png",   5.709, 3.179),  # slide 3, lower-right figure
    (5, "dashboard.png", 0.197, 1.160),  # slide 6, large left wireframe
]
TOL_EMU = Inches(1.0)  # allow up to 1in of manual nudging before we refuse


def pictures(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def swap(slide, png, exp_left, exp_top):
    tl, tt = Inches(exp_left), Inches(exp_top)
    pics = pictures(slide)
    if not pics:
        raise RuntimeError(f"no pictures on slide for {png}")
    best = min(pics, key=lambda s: (s.left - tl) ** 2 + (s.top - tt) ** 2)
    dist = ((best.left - tl) ** 2 + (best.top - tt) ** 2) ** 0.5
    if dist > TOL_EMU:
        raise RuntimeError(
            f"closest picture for {png} is {dist/914400:.2f}in from expected "
            f"position - refusing (image may have been moved manually).")
    rId = best._element.blip_rId
    image_part = best.part.related_part(rId)
    with open(os.path.join(A, png), "rb") as f:
        image_part._blob = f.read()
    print(f"swapped {png}: matched picture at "
          f"({best.left/914400:.3f}in, {best.top/914400:.3f}in), "
          f"offset {dist/914400:.3f}in, "
          f"size {best.width/914400:.3f}x{best.height/914400:.3f}in")


def main():
    prs = Presentation(OUT)
    for idx, png, l, t in TARGETS:
        swap(prs.slides[idx], png, l, t)
    prs.save(OUT)
    print("SAVED:", OUT)


if __name__ == "__main__":
    main()
