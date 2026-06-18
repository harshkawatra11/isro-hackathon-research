# -*- coding: utf-8 -*-
"""Lay out the SuryaSetu PS-15 pitch into ISRO's 10-slide template (chrome preserved).
Text slides use native light-theme typography; only genuine data visualisations are images.
All text uses plain keyboard characters only."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

SRC = "[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
OUT = "SuryaSetu - ISRO BAH 2026 - PS15 - Idea Submission.pptx"
A   = "render_assets"

INK    = RGBColor(0x16, 0x1d, 0x2e)
ACCENT = RGBColor(0xc2, 0x56, 0x0a)
SOFTB  = RGBColor(0x1f, 0x6f, 0xb8)
MUTED  = RGBColor(0x51, 0x5b, 0x73)
GREEN  = RGBColor(0x17, 0x82, 0x55)
AMBER  = RGBColor(0xb0, 0x7a, 0x00)
PURPLE = RGBColor(0x6d, 0x47, 0xc9)
RED    = RGBColor(0xc0, 0x2a, 0x42)
WHITE  = RGBColor(0xff, 0xff, 0xff)
NAVY   = RGBColor(0x0b, 0x10, 0x1d)
CARD   = RGBColor(0xf4, 0xf6, 0xfb)
CARD2  = RGBColor(0xeb, 0xef, 0xf8)
BORDER = RGBColor(0xd6, 0xdd, 0xec)
BODY   = RGBColor(0x12, 0x16, 0x20)   # near-black for readable body text
LINKB  = RGBColor(0x6d, 0xb4, 0xff)   # link blue on navy banners
F = "Segoe UI"
TITLE_TOP = 0.70

p = Presentation(SRC)
S = p.slides

# clear ISRO's grey instruction placeholders on content slides 3-9 (keep the branded PICTURE)
for idx in range(2, 9):
    for sh in S[idx].shapes:
        if sh.has_text_frame:
            sh.text_frame.clear()

# ---------------- helpers ----------------
def _set(run, text, size, color, bold=False, italic=False, font=F):
    run.text = text; run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color; run.font.name = font

def title_bar(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.34), Inches(TITLE_TOP), Inches(9.3), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    _set(tf.paragraphs[0].add_run(), text, 19, ACCENT, bold=True)
    return tb

def lines(slide, items, left, top, width, height=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height or 3.0))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for text, size, color, bold, italic, bullet in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        para.space_after = Pt(3)
        _set(para.add_run(), ("- " if bullet else "") + text, size, color, bold, italic)
    return tb

def caption(slide, text, left, top, width, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, bold=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    para = tf.paragraphs[0]; para.alignment = align
    _set(para.add_run(), text, size, color, bold=bold)
    return tb

def banner(slide, text, left, top, width, height, fill=NAVY, tc=WHITE, size=11, bold=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = ACCENT; sh.line.width = Pt(1.25)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    _set(para.add_run(), text, size, tc, bold=bold)
    return sh

def img(slide, name, left, top, max_w, max_h, align="center"):
    path = os.path.join(A, name)
    iw, ih = Image.open(path).size; ar = iw/ih
    w = max_w; h = w/ar
    if h > max_h:
        h = max_h; w = h*ar
    if align == "center":
        left = left + (max_w - w)/2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return w, h

def feature_item(slide, x, y, colw, num, color, title, body):
    # solid square number marker (sharp corners)
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.42), Inches(0.42))
    sq.fill.solid(); sq.fill.fore_color.rgb = color; sq.line.fill.background(); sq.shadow.inherit = False
    st = sq.text_frame
    st.margin_left = 0; st.margin_right = 0; st.margin_top = 0; st.margin_bottom = 0
    st.vertical_anchor = MSO_ANCHOR.MIDDLE
    sp = st.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    _set(sp.add_run(), num, 13, WHITE, bold=True)
    tb = slide.shapes.add_textbox(Inches(x+0.56), Inches(y-0.05), Inches(colw-0.6), Inches(0.34))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    _set(tf.paragraphs[0].add_run(), title, 12.5, INK, bold=True)
    bb = slide.shapes.add_textbox(Inches(x+0.56), Inches(y+0.30), Inches(colw-0.6), Inches(0.72))
    bf = bb.text_frame; bf.word_wrap = True; bf.auto_size = MSO_AUTO_SIZE.NONE
    _set(bf.paragraphs[0].add_run(), body, 9.7, BODY)
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y+1.06),
                                    Inches(x+colw-0.2), Inches(y+1.06))
    cn.line.color.rgb = BORDER; cn.line.width = Pt(0.75)

def rich(slide, paras, left, top, width, height, size=9, color=BODY, align=PP_ALIGN.LEFT,
         gap=3, fill=None, tc=None):
    """Multi-run paragraphs. Each para is a list of runs; a run is a tuple
    (text, bold=False, italic=False, size=None, color=None, link=None).
    If fill is given, draws a filled rectangle (banner); else a plain textbox."""
    if fill is not None:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = ACCENT; sh.line.width = Pt(1.25)
        tf = sh.text_frame
        tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if color is BODY: color = tc or WHITE
    else:
        sh = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = sh.text_frame
    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for para in paras:
        para_p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        para_p.alignment = align; para_p.space_after = Pt(gap)
        for run in para:
            text = run[0]
            bold = run[1] if len(run) > 1 else False
            ital = run[2] if len(run) > 2 else False
            rs   = run[3] if len(run) > 3 and run[3] else size
            rc   = run[4] if len(run) > 4 and run[4] else color
            r = para_p.add_run(); _set(r, text, rs, rc, bold, ital)
            if len(run) > 5 and run[5]:
                r.hyperlink.address = run[5]
            if len(run) > 6 and run[6]:
                r.font.underline = True
    return sh

# =================================================== SLIDE 1 - Title
for sh in S[0].shapes:
    if not sh.has_text_frame: continue
    t = sh.text_frame.text
    if t.startswith("Team Name"):       val = "  SuryaSetu"
    elif t.startswith("Problem Statement"): val = "  PS-15 | Nowcasting and Forecasting of Solar Flares from Aditya-L1 X-ray data"
    elif t.startswith("Team Leader"):    val = "  Harsh Kawatra"
    else: continue
    _set(sh.text_frame.paragraphs[0].add_run(), val, 16, ACCENT, bold=True)
caption(S[0], "SoLEXS 2-22 keV     HEL1OS 8-150 keV     LightGBM + MiniROCKET     TSS-driven     Zero cost on a laptop",
        0.34, 5.18, 9.3, size=9.5, color=SOFTB, align=PP_ALIGN.LEFT, bold=True)

# =================================================== SLIDE 2 - Team
for sh in S[1].shapes:
    if sh.has_table:
        tbl = sh.table
        cells = {
            (0,0): ("Team Leader", "Harsh Kawatra", "Delhi Technological University (DTU)"),
            (0,1): ("Team Member-1", "Gursimran Kaur", "Guru Gobind Singh Indraprastha University (GGSIPU)"),
            (1,0): ("Team Member-2", "-", "-"),
            (1,1): ("Team Member-3", "-", "-"),
        }
        for (ri,ci),(role,name,college) in cells.items():
            tf = tbl.cell(ri,ci).text_frame; tf.clear(); tf.word_wrap = True
            _set(tf.paragraphs[0].add_run(), role, 13, ACCENT, bold=True)
            for label,val in [("Name: ",name),("College: ",college)]:
                par = tf.add_paragraph()
                _set(par.add_run(), label, 11, INK, bold=True)
                _set(par.add_run(), val, 11, INK)
rich(S[1], [[
    ("A two-person, two-institute team: ", False, False, 10.5, BODY),
    ("full-stack and AI/ML engineering", True, False, 10.5, BODY),
    (" applied to ISRO's own solar mission.", False, True, 10.5, BODY),
]], 0.5, 5.15, 9.0, 0.4, align=PP_ALIGN.CENTER)

# =================================================== SLIDE 3 - Opportunity & USP
title_bar(S[2], "Opportunity and USP: why SuryaSetu wins")
rich(S[2], [
    [('Most flare tools use a single channel from Earth orbit, and most "forecasts" are really just detections after the flare has already begun. ', False, False, 9.5),
     ("SuryaSetu is built differently on three axes:", True, False, 9.5)],
    [("1. Dual-payload fusion.", True, False, 9.5, ACCENT)],
    [("SoLEXS measures ", ), ("soft X-rays", True), (" (the heat a flare gives off) and HEL1OS measures ", ), ("hard X-rays", True), (" (the instant energy is released).", )],
    [("We merge both onto a ", ), ("single timeline", True), (" from the L1 vantage, so the model reads the ", ), ("full physics", True), (" of a flare instead of half of it - ", ), ("the biggest reason our forecasts beat single-channel tools.", False, True)],
    [("2. Physics-grounded lead time.", True, False, 9.5, SOFTB)],
    [("Hard X-rays surge several minutes before the damaging ", ), ("soft-X-ray peak", True), (", a real published relationship called the ", ), ("Neupert effect", True), (".", )],
    [("We convert that natural head-start into the actual ", ), ("warning time", True), (" an operator needs, so the forecast rests on solar physics, ", ), ("not a curve fit that only matches past data.", False, True)],
    [("3. Honest, operational evaluation.", True, False, 9.5, GREEN)],
    [('Flares are rare, so a model that always says "no flare" looks ', ), ("99% accurate", True), (" and is useless.", )],
    [("We score with the ", ), ("True Skill Statistic", True), (", calibrate every probability, and validate only on ", ), ("future time blocks", True), (" - the standard ISRO and the space-weather community actually trust.", )],
], 0.34, 1.08, 4.58, height=4.05, size=9, color=BODY, gap=2)
img(S[2], "spectrum.png",  5.02, 1.18, 4.68, 1.85, align="center")
img(S[2], "neupert.png",   5.02, 3.10, 4.68, 1.95, align="center")
banner(S[2], "USP: the only dual-payload, L1-vantage system that turns the Neupert effect into calibrated, TSS-validated lead-time warnings, reproducible at zero cost on a laptop.",
       0.34, 5.18, 9.32, 0.34, fill=NAVY, tc=WHITE, size=9.5)

# =================================================== SLIDE 4 - Features (native light)
title_bar(S[3], "Features offered by the solution")
feats = [
    ("01", ACCENT, "Automated nowcast catalogue", "Real-time soft and hard detection merged into one master catalogue: onset, peak, end, class and SNR."),
    ("02", SOFTB,  "Calibrated forecast + lead time", "P(flare within N minutes) with calibrated confidence and a measured lead-time distribution."),
    ("03", GREEN,  "Live dashboard + visual alerts", "Dual light curves, hard/soft ratio and a green, amber, red banner that flips before the peak."),
    ("04", AMBER,  "Full A-to-X dynamic range", "Adaptive baselining plus dual SoLEXS detectors catch faint A/B/C and giant M/X flares."),
    ("05", PURPLE, "Explainable alerts (SHAP)", "Every forecast states which precursor fired it, so scientists can trust and audit it."),
    ("06", RED,    "Offline and reproducible", "No cloud, no cost, no proprietary software: end-to-end on one laptop from public data."),
]
cols_x = [0.45, 5.15]; colw = 4.45; rows_y = [1.24, 2.42, 3.60]
for i, (num, col, ttl, body) in enumerate(feats):
    feature_item(S[3], cols_x[i % 2], rows_y[i // 2], colw, num, col, ttl, body)
rich(S[3], [
    [("How it works in practice: ", True, False, 9, AMBER),
     ("the system reads both Aditya-L1 X-ray channels continuously, ", ),
     ("automatically logs every flare", True), (" it detects into the catalogue, and ", ),
     ("forecasts the next flare's strength and timing", True), (" with a ", ),
     ("calibrated confidence", True), (".", )],
    [("The operator watches one dashboard; when risk crosses a tuned threshold the ", ),
     ("banner turns red", True, False, 9, AMBER), (" with the ", ),
     ("minutes of lead time", True), (" remaining, and ", ),
     ("SHAP", True), (" names the exact precursor that triggered it.", )],
], 0.40, 4.78, 9.2, 0.70, fill=NAVY, color=WHITE, size=9, align=PP_ALIGN.CENTER, gap=2)

# =================================================== SLIDE 5 - Process flow / use-case
title_bar(S[4], "Process flow and use-case")
img(S[4], "process_flow.png", 0.28, 1.22, 6.25, 3.40, align="center")
img(S[4], "usecase.png",      6.62, 1.22, 3.10, 3.40, align="center")
rich(S[4], [
    [("In plain terms: ", True, False, 8.4, AMBER), ("raw satellite files are ", ), ("cleaned", True), (" (bad seconds removed) and merged into one ", ), ("1-second timeline", True), (".", )],
    [("That stream splits into ", ), ("two engines", True), (" running side by side: a ", ), ("classical detector", True), (" that logs flares the moment they start, and an ", ), ("ML forecaster", True), (" that predicts the next one.", )],
    [("Both feed the ", ), ("catalogue, the alerts, and the live dashboard", True), (".", )],
], 0.30, 4.64, 6.25, 0.95, size=8.4, color=BODY, gap=2)
rich(S[4], [
    [("Two users.", True)],
    [("An ", ), ("operations engineer", True), (" receives an alert and safe-modes satellites or holds GPS-critical work;", )],
    [("a ", ), ("researcher", True), (" browses the catalogue, replays any past day, and audits why each alert fired.", )],
], 6.62, 4.64, 3.10, 0.95, size=8.4, color=BODY, gap=2)

# =================================================== SLIDE 6 - Wireframe / dashboard
title_bar(S[5], "Wireframe: the operator dashboard")
img(S[5], "dashboard.png", 0.28, 1.16, 6.25, 3.60, align="left")
img(S[5], "real_lightcurve.png", 6.60, 1.16, 3.14, 1.05, align="center")
rich(S[5], [[("Renders ", False, True, 8, BODY), ("real Aditya-L1 data", True, False, 8, BODY), (" (2026-06-15)", False, True, 8, BODY)]],
     6.60, 2.16, 3.14, 0.3, align=PP_ALIGN.CENTER)
states = [
    ("QUIET", GREEN, [
        [("The Sun is ", ), ("calm", True), (" and ", ), ("no action is needed", True), (".", )],
        [("The system keeps logging the ", ), ("quiet-Sun baseline", True), (" it measures against.", )],
    ]),
    ("WATCH", AMBER, [
        [("The ", ), ("hard/soft ratio is climbing", True), (", the physics precursor that often comes before a flare.", )],
        [("The operator is ", ), ("put on notice", True), (" before anything is certain.", )],
    ]),
    ("ALERT", RED, [
        [("A flare is forecast.", True)],
        [("The banner shows the predicted ", ), ("class", True), (", the ", ), ("minutes of lead time", True), (", and a ", ), ("calibrated confidence", True), (" so the operator can act.", )],
    ]),
]
y = 2.48
for label, col, descp in states:
    dotsh = S[5].shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.62), Inches(y+0.03), Inches(0.14), Inches(0.14))
    dotsh.fill.solid(); dotsh.fill.fore_color.rgb = col; dotsh.line.fill.background(); dotsh.shadow.inherit = False
    tb = S[5].shapes.add_textbox(Inches(6.82), Inches(y-0.04), Inches(2.92), Inches(0.86))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    _set(tf.paragraphs[0].add_run(), label, 10.5, col, bold=True)
    for para in descp:
        pp = tf.add_paragraph(); pp.space_after = Pt(1)
        for run in para:
            text = run[0]; bold = run[1] if len(run) > 1 else False
            _set(pp.add_run(), text, 8, BODY, bold)
    y += 0.92
rich(S[5], [
    [("One screen for the operator: ", True), ("the live soft and hard light curves with a ", ), ("moving playhead", True), (", the rising ", ), ("hard/soft ratio", True), (" that acts as the early warning, a ", ), ("probability gauge", True), (" for the next N minutes, and the running ", ), ("flare catalogue", True), (".", )],
    [("The banner ", ), ("flips red", True), (" before the ", ), ("soft-X-ray peak", True), (", and ", ), ("that gap is the warning time we deliver.", False, True)],
], 0.30, 4.80, 6.25, 0.70, size=8.2, color=BODY, gap=2)

# =================================================== SLIDE 7 - Architecture
title_bar(S[6], "Architecture of the proposed solution")
img(S[6], "architecture.png", 0.28, 1.10, 6.05, 3.18, align="left")
img(S[6], "tss.png",                6.38, 1.10, 3.34, 2.02, align="center")
img(S[6], "feature_importance.png", 6.38, 3.22, 3.34, 1.78, align="center")
rich(S[6], [
    [("How it works: ", True, False, 8, AMBER), ("the cleaned soft and hard light curves feed ", ), ("two paths at once", True), (".", )],
    [("A ", ), ("classical signal-processing detector", True), (" (adaptive baseline plus rise-rate) handles nowcasting, ", ), ("interpretable and training-free", False, True), (".", )],
    [("In parallel, a ", ), ("LightGBM", True), (" gradient-boosted model on Neupert-physics features and a ", ), ("MiniROCKET", True), (" convolutional time-series model on raw windows are fused and ", ), ("isotonic-calibrated", True), (" into one probability, ", ), ("P(flare within N minutes)", True), (".", )],
    [("Every result is validated with ", ), ("time-blocked cross-validation", True), (" on the ", ), ("True Skill Statistic", True), (", and ", ), ("SHAP", True), (" attributes each alert to the precursor that drove it.", )],
], 0.28, 4.34, 6.05, 1.15, size=8, color=BODY, gap=1)

# =================================================== SLIDE 8 - Technologies (native light table)
title_bar(S[7], "Technologies used in the solution")
stack = [
    ("FITS I/O",          "astropy",             "standard reader for SoLEXS / HEL1OS products + GTI", "BSD",    ACCENT),
    ("Numerics",          "numpy, pandas",       "time-series alignment, resampling, fusion",          "BSD",    ACCENT),
    ("Signal",            "scipy",               "detrend, smooth, derivative, adaptive baseline",     "BSD",    ACCENT),
    ("Forecaster A",      "LightGBM",            "fast, interpretable, sample-efficient on rare events","MIT",   SOFTB),
    ("Forecaster B",      "MiniROCKET (sktime)", "CPU, deterministic, strong on raw time series",      "BSD",    SOFTB),
    ("Calibrate, metrics","scikit-learn",        "isotonic calibration, TSS/HSS/ROC/Brier, blocked CV","BSD",    SOFTB),
    ("Explainability",    "SHAP",                "per-alert precursor attribution (the why)",          "MIT",    PURPLE),
    ("Figures",           "matplotlib, plotly",  "light curves, reliability, lead-time plots",         "BSD",    GREEN),
    ("Interface",         "Streamlit",           "fastest path to dual curves, alerts and replay",     "Apache", GREEN),
    ("Environment",       "Python 3.11 (uv)",    "prebuilt wheels, avoids the 3.14 astropy build break","PSF",   MUTED),
]
nr = len(stack) + 1
tw, th = 7.55, 3.62
tbl = S[7].shapes.add_table(nr, 4, Inches(0.34), Inches(1.08), Inches(tw), Inches(th)).table
tbl.first_row = False; tbl.horz_banding = False
widths = [1.4, 1.62, 3.72, 0.81]
for ci, wv in enumerate(widths): tbl.columns[ci].width = Inches(wv)
heads = ["Layer", "Tool", "Why this tool", "License"]
for ci, htext in enumerate(heads):
    c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    par = c.text_frame.paragraphs[0]; _set(par.add_run(), htext, 9.5, WHITE, bold=True)
for ri, (layer, tool, why, lic, col) in enumerate(stack, start=1):
    rowfill = CARD if ri % 2 else CARD2
    for ci, (val, color, bold, fam) in enumerate([
            (layer, col, True, F), (tool, INK, True, "Consolas"),
            (why, BODY, False, F), (lic, INK, False, "Consolas")]):
        c = tbl.cell(ri, ci); c.fill.solid(); c.fill.fore_color.rgb = rowfill
        c.margin_top = Inches(0.015); c.margin_bottom = Inches(0.015); c.margin_left = Inches(0.08)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        par = c.text_frame.paragraphs[0]
        _set(par.add_run(), val, 9, color, bold=bold, font=fam)
lines(S[7], [
    ("Every layer, justified.", 11, ACCENT, True, False, False),
    ("", 5, INK, False, False, False),
    ("The prototype runs on a modest laptop (i5 / RTX 3050 / 16 GB) - proof that the method, not the hardware, does the work.", 9, INK, False, False, False),
    ("", 4, INK, False, False, False),
    ("The same open-source stack scales straight onto ISRO's far larger infrastructure for full-mission throughput.", 9, INK, False, False, False),
    ("", 4, INK, False, False, False),
    ("In ISRO's own tradition: the best results from the least resources.", 9, GREEN, True, True, False),
], 8.05, 1.20, 1.78, height=3.55)
VURL = "https://isro-hack-research.vercel.app"
rich(S[7], [
    [("Proof of work and open research: ", True, False, 9, WHITE),
     ("isro-hack-research.vercel.app", True, False, 9, LINKB, VURL, True),
     ("  [ Click to Know ]", True, False, 9, LINKB, VURL, True),
     (" - a complete, first-principles companion to the science, data, and methodology behind SuryaSetu,", False, False, 8.7, WHITE)],
    [("written and published so any engineer can understand and rebuild the system.", False, True, 8.7, WHITE)],
], 0.34, 4.82, 9.32, 0.62, fill=NAVY, color=WHITE, size=9, align=PP_ALIGN.CENTER, gap=1)

# =================================================== SLIDE 9 - Cost
title_bar(S[8], "Estimated implementation cost")
lines(S[8], [
    ("Prototype (our submission): effectively zero rupees. It needs only an internet connection, a capable GPU (enough VRAM) and a good CPU, all of which we already own.", 10.5, ACCENT, True, False, False),
], 0.34, 1.12, 9.3, height=0.55)
rows = [
    ("Item","Cost","Note"),
    ("Aditya-L1 SoLEXS + HEL1OS data","Rs 0","free, ISSDC PRADAN (registered)"),
    ("GOES flare catalog (labels)","Rs 0","free, NOAA SWPC"),
    ("Software stack","Rs 0","all open-source (BSD/MIT/Apache)"),
    ("Compute","Rs 0","existing laptop, i5 / RTX 3050 / 16 GB"),
    ("Storage","Rs 0","existing 256 GB drive (~5 GB used)"),
    ("Hosting (demo)","Rs 0","Streamlit local / free HF Spaces"),
    ("TOTAL (prototype)","Rs 0","only real input is engineering time"),
]
gt = S[8].shapes.add_table(len(rows), 3, Inches(0.34), Inches(1.62), Inches(5.55), Inches(2.7)).table
gt.first_row = False; gt.horz_banding = False
gt.columns[0].width = Inches(2.7); gt.columns[1].width = Inches(0.7); gt.columns[2].width = Inches(2.15)
for ri,row in enumerate(rows):
    head = ri==0; tot = ri==len(rows)-1
    for ci,val in enumerate(row):
        c = gt.cell(ri,ci); c.fill.solid()
        c.fill.fore_color.rgb = NAVY if head else (RGBColor(0xfb,0xf1,0xe6) if tot else (CARD if ri%2 else CARD2))
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02); c.margin_left = Inches(0.08)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        par = c.text_frame.paragraphs[0]
        _set(par.add_run(), val, 9.5, WHITE if head else (ACCENT if tot else INK), bold=head or tot)
img(S[8], "data_funnel.png", 6.05, 1.62, 3.65, 2.78, align="center")
rich(S[8], [
    [("Cost at ISRO scale", True, False, 9.5, AMBER), (" (operational, scoped on selection):", False, True, 9, WHITE)],
    [("compute at scale", True), ("  (GPU cluster or cloud)   |   ", ), ("AI tooling and API keys", True), ("   |   ", ), ("redundant full-mission storage", True), ("   |   ", ), ("internet and data egress", True), ("   |   ", ), ("engineering and MLOps staff", True), (".", )],
], 0.34, 4.70, 9.32, 0.72, fill=NAVY, color=WHITE, size=9, align=PP_ALIGN.CENTER, gap=2)

# =================================================== SLIDE 10 - Closing
# left as the ISRO template's own closing slide (no added banner or references line)

p.save(OUT)
print("SAVED:", OUT)
