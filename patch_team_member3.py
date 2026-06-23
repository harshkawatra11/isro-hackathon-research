# -*- coding: utf-8 -*-
"""Surgically add Team Member-3 (Dayita Arora) into cell (1,1) of the team table
on slide 2 of the live deck, and update the caption from three-person to four-person.
No other shapes, text, or images are touched."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

OUT = "SuryaSetu - ISRO BAH 2026 - PS15 - Idea Submission.pptx"

INK   = RGBColor(0x16, 0x1d, 0x2e)
ACCENT = RGBColor(0xc2, 0x56, 0x0a)
F = "Segoe UI"

def _set(run, text, size, color, bold=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = F

def patch_table(slide):
    for sh in slide.shapes:
        if not sh.has_table:
            continue
        tbl = sh.table
        tf = tbl.cell(1, 1).text_frame
        tf.clear(); tf.word_wrap = True
        _set(tf.paragraphs[0].add_run(), "Team Member-3", 13, ACCENT, bold=True)
        for label, val in [("Name: ", "Dayita Arora"),
                           ("College: ", "Ramjas College (DU)"),
                           ("Major: ", "B.Sc in Statistics")]:
            par = tf.add_paragraph()
            _set(par.add_run(), label, 11, INK, bold=True)
            _set(par.add_run(), val, 11, INK)
        print("patched table cell (1,1) with Dayita Arora")
        return True
    return False

def patch_caption(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if "three-person" in run.text:
                    run.text = run.text.replace("three-person", "four-person")
                    print("patched caption: three-person -> four-person")
                    return True
                if "A three-person" in run.text:
                    run.text = run.text.replace("A three-person", "A four-person")
                    print("patched caption: A three-person -> A four-person")
                    return True
    return False

def main():
    prs = Presentation(OUT)
    slide2 = prs.slides[1]
    ok_table = patch_table(slide2)
    ok_caption = patch_caption(slide2)
    if not ok_table:
        raise RuntimeError("Could not find team table on slide 2")
    if not ok_caption:
        print("warning: caption text not found — may already be updated or in a different run")
    prs.save(OUT)
    print("SAVED:", OUT)

if __name__ == "__main__":
    main()
